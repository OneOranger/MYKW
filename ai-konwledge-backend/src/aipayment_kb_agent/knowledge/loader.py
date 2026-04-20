from __future__ import annotations

import logging
from pathlib import Path

import pandas as pd
from docx import Document as DocxDocument
from pypdf import PdfReader
from pptx import Presentation

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {
    ".txt",
    ".md",
    ".pdf",
    ".docx",
    ".pptx",
    ".csv",
    ".xlsx",
}


def gather_supported_files(target: Path) -> list[Path]:
    if target.is_file():
        return [target] if target.suffix.lower() in SUPPORTED_EXTENSIONS else []
    files = []
    for path in target.rglob("*"):
        if path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS:
            files.append(path)
    return files


def read_document(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in {".txt", ".md"}:
        return path.read_text(encoding="utf-8", errors="ignore")
    if suffix == ".pdf":
        reader = PdfReader(str(path))
        return "\n".join((page.extract_text() or "") for page in reader.pages)
    if suffix == ".docx":
        doc = DocxDocument(str(path))
        return "\n".join(paragraph.text for paragraph in doc.paragraphs)
    if suffix == ".pptx":
        deck = Presentation(str(path))
        blocks: list[str] = []
        for slide in deck.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    blocks.append(shape.text)
        return "\n".join(blocks)
    if suffix in {".csv", ".xlsx"}:
        df = pd.read_csv(path) if suffix == ".csv" else pd.read_excel(path)
        return df.to_csv(index=False)
    logger.warning("Unsupported file skipped: %s", path)
    return ""
